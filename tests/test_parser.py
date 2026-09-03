"""
Unit tests for DocumentParser and Markdown section chunking.
"""

from rag_agent.parser import DocumentParser, chunk_markdown


def test_chunk_markdown_preserves_tables_and_headers():
    markdown_text = """
# Product Specifications

Here is the introduction to TechPro X1.

| Model | RAM | Storage | Price |
| --- | --- | --- | --- |
| TechPro X1 | 16GB | 512GB | $999 |
| TechPro Ultra | 32GB | 1TB | $1499 |

## Support Guidelines
Contact customer support for pass reset.
"""
    chunks = chunk_markdown(markdown_text)

    # Verify tables and headers are parsed as distinct intact chunks
    assert len(chunks) >= 2
    # Ensure table block remains contiguous
    table_chunk = [c for c in chunks if "|" in c][0]
    assert "| Model | RAM | Storage | Price |" in table_chunk
    assert "| TechPro Ultra | 32GB | 1TB | $1499 |" in table_chunk


def test_document_parser_txt_file():
    txt_content = b"First section.\n\nSecond section with details."
    chunks, engine = DocumentParser.parse_file(txt_content, "sample.txt")

    assert len(chunks) == 2
    assert chunks[0] == "First section."
    assert "Text Parser" in engine


def test_document_parser_docx_fallback():
    import io
    import docx
    doc = docx.Document()
    doc.add_paragraph("Paragraph inside docx document.")
    table = doc.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"

    stream = io.BytesIO()
    doc.save(stream)
    docx_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(docx_bytes, "sample.docx")
    assert len(chunks) > 0
    assert "docx" in engine.lower() or "liteparse" in engine.lower()


def test_document_parser_xlsx_fallback():
    import io
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "DataSheet"
    ws.append(["Item", "Price"])
    ws.append(["Laptop", "$999"])

    stream = io.BytesIO()
    wb.save(stream)
    xlsx_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(xlsx_bytes, "data.xlsx")
    assert len(chunks) > 0
    assert "openpyxl" in engine.lower() or "liteparse" in engine.lower()


def test_document_parser_pptx_fallback():
    import io
    import pptx
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Presentation Title"

    stream = io.BytesIO()
    prs.save(stream)
    pptx_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(pptx_bytes, "deck.pptx")
    assert len(chunks) > 0
    assert "pptx" in engine.lower() or "liteparse" in engine.lower()


def test_document_parser_image_info():
    import io
    from PIL import Image
    img = Image.new("RGB", (100, 100), color="blue")
    stream = io.BytesIO()
    img.save(stream, format="PNG")
    img_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(img_bytes, "logo.png")
    assert len(chunks) > 0
    assert any(k in engine.lower() for k in ["pillow", "liteparse", "ocr", "image", "vision"])


def test_document_parser_vision_llm_mock(monkeypatch):
    """Vision-OCR path. Groq currently serves no vision model, so the model list
    is empty by default; inject one to keep the code path covered."""
    import io
    from unittest.mock import MagicMock
    from PIL import Image

    monkeypatch.setattr("rag_agent.parser.vision_models_for", lambda: ["test-vision-model"])

    img = Image.new("RGB", (200, 200), color="white")
    stream = io.BytesIO()
    img.save(stream, format="PNG")
    img_bytes = stream.getvalue()

    mock_groq_client = MagicMock()
    mock_choice = MagicMock()
    mock_choice.message.content = "## Extracted Text:\nProduct Specs: TechPro Laptop\nPrice: $999"
    mock_response = MagicMock()
    mock_response.choices = [mock_choice]
    mock_groq_client.chat.completions.create.return_value = mock_response

    chunks, engine = DocumentParser.parse_file(img_bytes, "specs.png", groq_client=mock_groq_client)
    assert len(chunks) > 0
    assert "TechPro Laptop" in "\n".join(chunks)
    assert "Vision LLM Engine" in engine


def _image_only_pdf() -> bytes:
    """A PDF with no text layer (a stand-in for a scanned document)."""
    import io
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (1240, 1754), "white")
    d = ImageDraw.Draw(img)
    for i in range(20):
        d.text((90, 90 + i * 60), f"Scanned line {i}: total equity 10,208,747", fill="black")
    buf = io.BytesIO()
    img.save(buf, format="PDF", resolution=150)
    return buf.getvalue()


def _text_native_pdf(lines: int = 12) -> bytes:
    """A minimal hand-built PDF carrying a real text layer.

    Written by hand because no PDF *writer* library (reportlab/fpdf) is a
    project dependency, and the positive case is what the OCR-skip decision
    hinges on.
    """
    text_ops = "BT /F1 12 Tf 72 720 Td 14 TL\n" + "".join(
        f"(Total equity was 10,208,747 in 2019 - line {i}) Tj T*\n" for i in range(lines)
    ) + "ET"
    stream = f"<< /Length {len(text_ops)} >>\nstream\n{text_ops}\nendstream"
    objects = [
        "<< /Type /Catalog /Pages 2 0 R >>",
        "<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        "<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        "/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        stream,
        "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = "%PDF-1.4\n"
    offsets = []
    for i, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n{body}\nendobj\n"
    xref_at = len(out)
    out += f"xref\n0 {len(objects)+1}\n0000000000 65535 f \n"
    out += "".join(f"{off:010d} 00000 n \n" for off in offsets)
    out += (f"trailer\n<< /Size {len(objects)+1} /Root 1 0 R >>\n"
            f"startxref\n{xref_at}\n%%EOF\n")
    return out.encode("latin-1")


def test_has_text_layer_true_for_text_native_pdf():
    """Text-native PDFs must skip OCR - that decision is a ~300x speed difference."""
    assert DocumentParser.has_text_layer(_text_native_pdf()) is True


def test_has_text_layer_false_for_image_only_pdf():
    """Scanned-style PDFs must still be routed through OCR."""
    assert DocumentParser.has_text_layer(_image_only_pdf()) is False


def test_has_text_layer_false_for_non_pdf_bytes():
    assert DocumentParser.has_text_layer(b"this is not a pdf at all") is False


def test_document_parser_image_skips_vision_when_no_model_configured():
    """With no vision model available, image parsing falls through to LiteParse OCR."""
    import io
    from unittest.mock import MagicMock
    from PIL import Image

    img = Image.new("RGB", (200, 200), color="white")
    stream = io.BytesIO()
    img.save(stream, format="PNG")
    img_bytes = stream.getvalue()

    mock_groq_client = MagicMock()
    chunks, engine = DocumentParser.parse_file(img_bytes, "blank.png", groq_client=mock_groq_client)

    # No vision model configured -> the client is never called for OCR
    mock_groq_client.chat.completions.create.assert_not_called()
    assert "Vision LLM Engine" not in engine


def test_document_parser_text_image_ocr():
    import io
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (400, 100), color=(255, 255, 255))
    d = ImageDraw.Draw(img)
    d.text((10, 10), "TechPro Laptop OCR Test", fill=(0, 0, 0))

    stream = io.BytesIO()
    img.save(stream, format="PNG")
    img_bytes = stream.getvalue()

    chunks, engine = DocumentParser.parse_file(img_bytes, "test_ocr.png")
    assert len(chunks) > 0


